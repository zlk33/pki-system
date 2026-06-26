#!/usr/bin/env python3
"""
Generuje prezentacje PPTX projektu Studenckie PKI.

Wynik: docs/generated/studenckie-pki-prezentacja.pptx

Uzycie:
  pip install -r scripts/requirements-docs.txt
  python scripts/generate_presentation.py
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUTPUT = ROOT / "docs" / "generated" / "studenckie-pki-prezentacja.pptx"

TEAM_DISPLAY = [
    "Jakub Zalewski",
    "Hubert Rachwalik",
    "Jakub Zaborski",
    "Mateusz Rutkowski",
    "Michał Dęby",
    "Robert Graff",
    "Wojciech Bojko",
]

ACCENT = RGBColor(0x1D, 0x4E, 0xD8)
TEXT = RGBColor(0x11, 0x18, 0x27)
MUTED = RGBColor(0x4B, 0x55, 0x63)
PLACEHOLDER_BG = RGBColor(0xF3, 0xF4, 0xF6)
PLACEHOLDER_BORDER = RGBColor(0x9C, 0xA3, 0xAF)

SLIDES = [
    {
        "title": "Studenckie PKI",
        "subtitle": "System zarządzania infrastrukturą klucza publicznego",
        "bullets": [
            "Zespół projektowy:",
            *TEAM_DISPLAY,
            "",
            "Stack: Python (FastAPI) + PostgreSQL + Next.js + Docker",
        ],
        "notes": "Przedstawiamy projekt edukacyjnego PKI — od Root CA, przez certyfikaty, po podpisy i weryfikację dokumentów.",
        "image_hint": "Zdjęcie zespołu lub logo projektu",
        "layout": "title",
    },
    {
        "title": "Cel projektu",
        "bullets": [
            "Działający, lokalny system PKI z panelem webowym",
            "Rejestracja, logowanie, certyfikaty użytkowników",
            "Podpisywanie dokumentów i weryfikacja podpisu",
            "Panel admina: użytkownicy, CRL, dziennik zdarzeń",
        ],
        "notes": "To nie jest produkcyjne HSM-PKI, ale kompletny przepływ X.509, CSR, CRL i podpisów odłączonych.",
        "image_hint": "Zrzut ekranu — dashboard",
    },
    {
        "title": "Architektura",
        "bullets": [
            "Frontend Next.js (port 3000) — panel użytkownika",
            "Backend FastAPI (port 8000) — logika PKI w Pythonie",
            "PostgreSQL — użytkownicy, certyfikaty, audyt",
            "Docker Compose — 3 kontenery + wolumen /app/certs",
        ],
        "notes": "Cała kryptografia jest w Pythonie. Frontend tylko zbiera dane i wyświetla wyniki.",
        "image_hint": "Diagram architektury lub docker compose",
    },
    {
        "title": "Backend (Python)",
        "bullets": [
            "main.py — REST API, JWT, orchestracja",
            "ca_manager.py — generowanie Root CA",
            "cert_manager.py — podpisywanie CSR",
            "crl_manager.py — lista CRL (PEM)",
            "audit.py — dziennik zdarzeń",
            "db.py — modele SQLAlchemy",
        ],
        "notes": "Backend jest sercem projektu — spełnia wymaganie przedmiotu dotyczące Pythona.",
        "image_hint": "Struktura katalogu backend/app",
    },
    {
        "title": "Baza danych",
        "bullets": [
            "users — konta, role admin / user",
            "certificates — PEM, status, unieważnienie",
            "requests — żądania CSR",
            "signed_documents — metadane podpisów",
            "audit_logs — historia operacji",
        ],
        "notes": "Schemat w db/init.sql. Lekkie migracje przy starcie backendu.",
        "image_hint": "Schemat tabel lub pgAdmin",
    },
    {
        "title": "Uwierzytelnianie i role",
        "bullets": [
            "Rejestracja — pierwszy użytkownik = admin",
            "Logowanie — JWT Bearer",
            "Admin: CA, użytkownicy, certyfikaty, CRL, audyt",
            "User: własne certyfikaty, CSR, podpisywanie",
        ],
        "notes": "Każde żądanie API (poza login/register) wymaga tokenu.",
        "image_hint": "Ekran logowania / rejestracji",
    },
    {
        "title": "Root CA",
        "bullets": [
            "Admin → Root CA → Inicjalizuj (jednorazowo)",
            "RSA-2048, certyfikat samopodpisany (10 lat)",
            "Klucz w /app/certs/root_ca.key (wolumen Docker)",
            "Bez Root CA nie da się wystawiać certyfikatów",
        ],
        "notes": "Operacja jednorazowa — fundament całego PKI.",
        "image_hint": "Panel Root CA po inicjalizacji",
    },
    {
        "title": "Wystawianie certyfikatów",
        "bullets": [
            "Tryb A — CSR: wniosek → zatwierdzenie przez admina",
            "Tryb B — admin wystawia certyfikat dla użytkownika",
            "Klucz generowany na serwerze (tryb demo)",
            "Pobieranie certyfikatu PEM i klucza prywatnego",
        ],
        "notes": "Tryb B umożliwia podpisywanie dokumentów w panelu.",
        "image_hint": "Formularz wystawiania certyfikatu",
    },
    {
        "title": "Podpisywanie dokumentów",
        "bullets": [
            "Upload pliku → podpis RSA-SHA256",
            "Plik .pki.sig.json: hash, podpis, timestamp, cert PEM",
            "Podpis odłączony — oryginał zostaje u użytkownika",
            "Format: studenckie-pki-v2",
        ],
        "notes": "Nie modyfikujemy oryginału — dostarczamy plik podpisu do weryfikacji.",
        "image_hint": "Ekran podpisywania + pobrany plik .pki.sig.json",
    },
    {
        "title": "Weryfikacja podpisu",
        "bullets": [
            "Upload: oryginalny dokument + plik podpisu",
            "Sprawdzenie hasha SHA-256",
            "Weryfikacja podpisu kryptograficznego",
            "Kontrola ważności i statusu certyfikatu",
        ],
        "notes": "Można zweryfikować na innym komputerze — plik podpisu jest samowystarczalny.",
        "image_hint": "Ekran weryfikacji — wynik OK",
    },
    {
        "title": "CRL i unieważnianie",
        "bullets": [
            "Unieważnienie certyfikatu → status REVOKED + powód",
            "CRL (Certificate Revocation List) w formacie PEM",
            "Pobieranie CRL z panelu admina",
            "Weryfikacja odrzuca unieważniony certyfikat",
        ],
        "notes": "CRL to standardowy element PKI — lista numerów seryjnych revoked certów.",
        "image_hint": "Pobieranie CRL / lista unieważnionych",
    },
    {
        "title": "Zarządzanie użytkownikami",
        "bullets": [
            "Lista kont, tworzenie użytkownika",
            "Zmiana roli admin / user",
            "Reset hasła",
            "Usuwanie (z ochroną ostatniego admina)",
        ],
        "notes": "Pierwszy zarejestrowany użytkownik zostaje adminem.",
        "image_hint": "Panel Użytkownicy",
    },
    {
        "title": "Dziennik zdarzeń",
        "bullets": [
            "Kto wykonał — imię, email, rola (snapshot)",
            "Co — akcja i opis po polsku",
            "Na czym — cel, ID zasobu",
            "Metadane — np. zmiany przy edycji użytkownika",
        ],
        "notes": "Pełna historia działań w systemie — przydatne na obronie.",
        "image_hint": "Ekran dziennika zdarzeń (audit log)",
    },
    {
        "title": "Demo na żywo",
        "bullets": [
            "docker compose up --build",
            "Rejestracja admina → init Root CA",
            "Utworzenie użytkownika i wystawienie certyfikatu",
            "Podpisanie pliku → weryfikacja",
            "Unieważnienie → ponowna weryfikacja (błąd)",
            "CRL + dziennik zdarzeń",
        ],
        "notes": "Przejdź scenariusz spokojnie — to najważniejsza część prezentacji.",
        "image_hint": "Zrzut ekranu z demo / terminal Docker",
    },
    {
        "title": "Ograniczenia i rozwój",
        "bullets": [
            "Brak HSM — klucz Root CA w pliku",
            "Brak OCSP — tylko CRL",
            "Brak Intermediate CA",
            "Możliwe: PKCS#12, OCSP, TSA, HSM",
        ],
        "notes": "Projekt edukacyjny — wiemy, co byłoby potrzebne w produkcji.",
        "image_hint": "Opcjonalnie: schemat docelowej architektury",
    },
    {
        "title": "Podsumowanie",
        "bullets": [
            "Działające PKI w Dockerze",
            "Python jako rdzeń kryptograficzny",
            "Auth, role, panel admina",
            "Certyfikaty, CSR, CRL, podpisy, weryfikacja",
            "Audit log i dokumentacja",
            "",
            "Pytania?",
        ],
        "notes": "Dziękujemy za uwagę. Jesteśmy gotowi na pytania.",
        "image_hint": "Zdjęcie zespołu lub podziękowanie",
    },
]


def _set_run_font(run, size: int, bold: bool = False, color: RGBColor = TEXT) -> None:
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_title_bar(slide, title: str) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.4), Inches(0.08), Inches(9.2), Inches(0.45))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    _set_run_font(run, 24, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))


def _add_bullets(slide, bullets: list[str], left: float, top: float, width: float, height: float) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    first = True
    for line in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if not line:
            p.text = ""
            continue
        p.text = line
        p.level = 0
        p.space_after = Pt(6)
        for run in p.runs:
            _set_run_font(run, 16 if line.startswith("Stack") else 18)


def _add_image_placeholder(slide, left: float, top: float, width: float, height: float, hint: str) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PLACEHOLDER_BG
    shape.line.color.rgb = PLACEHOLDER_BORDER
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = "📷"
    _set_run_font(r1, 28, color=MUTED)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "Wstaw zdjęcie"
    _set_run_font(r2, 14, bold=True, color=MUTED)

    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = hint
    _set_run_font(r3, 11, color=MUTED)


def _add_title_slide(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(2.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.45), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = data["title"]
    _set_run_font(r, 40, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.6))
    tf2 = sub_box.text_frame
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = data.get("subtitle", "")
    _set_run_font(r2, 20, color=RGBColor(0xDB, 0xEA, 0xFE))

    _add_bullets(slide, data["bullets"], 0.5, 2.5, 4.5, 3.5)
    _add_image_placeholder(slide, 5.3, 2.5, 4.2, 4.2, data.get("image_hint", "Zdjęcie"))

    if data.get("notes"):
        slide.notes_slide.notes_text_frame.text = data["notes"]


def _add_content_slide(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, data["title"])
    _add_bullets(slide, data["bullets"], 0.45, 0.85, 5.4, 6.0)
    _add_image_placeholder(slide, 6.1, 0.85, 3.45, 5.8, data.get("image_hint", "Zrzut ekranu"))

    if data.get("notes"):
        slide.notes_slide.notes_text_frame.text = data["notes"]


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for i, data in enumerate(SLIDES):
        if data.get("layout") == "title" and i == 0:
            _add_title_slide(prs, data)
        else:
            _add_content_slide(prs, data)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    return OUTPUT


def main() -> int:
    path = build_presentation()
    print(f"PPTX: {path}")
    print("\nOtworz plik w PowerPoint i wstaw zdjęcia w szare pola.")
    print("Notatki do wygłoszenia sa w zakladce 'Notatki' pod kazdym slajdem.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
